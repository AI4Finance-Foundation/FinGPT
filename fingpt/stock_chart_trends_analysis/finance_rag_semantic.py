# finance_rag_semantic.py
"""
RAG data layer for financial documents (PDF/DOCX/TXT/MD/CSV/HTML/PPTX):
- Extract text (PDF via PdfExtractor; others via format-specific libs)
- Semantic chunking (heading-aware + similarity-based)
- Finance embeddings -> FAISS index
- Retrieval returns context + sources
- NO LLM CALLS HERE

Optional deps by type:
pip install python-docx beautifulsoup4 lxml python-pptx
Base deps:
pip install langchain langchain-community langchain-text-splitters sentence-transformers faiss-cpu tiktoken
"""

import io
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Union

import numpy as np

# All PDF work stays in this extractor
from fingpt.stock_chart_trends_analysis.PDFExtractor import PdfExtractor, SimplePdfValidator

# Optional extractors for non-PDF
try:
    from bs4 import BeautifulSoup  # HTML
except Exception:
    BeautifulSoup = None

try:
    from docx import Document as DocxDocument  # DOCX
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation  # PPTX
except Exception:
    Presentation = None

# LangChain bits (no LLMs here)
from langchain.schema import Document
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings

# ---------------------------
# Configuration (embeddings + index)
# ---------------------------
_raw = os.getenv("FIN_EMBED_MODEL")
EMBED_MODEL = (_raw or "").strip() or "sentence-transformers/all-MiniLM-L6-v2"
INDEX_DIR = Path("./faiss_fin_index")
TOP_K = int(os.getenv("TOP_K", "6"))

# Semantic chunking controls
MAX_CHARS = 3500
MIN_CHARS = 800
EMBED_BREAK_DELTA = 0.08
HARD_HEADINGS = [
    r"\bItem\s+7A?\b", r"\bRisk Factors\b", r"\bMD&A\b",
    r"\bQuantitative and Qualitative\b", r"\bConsolidated Statements\b",
    r"\bNotes to Consolidated\b", r"\bLiquidity and Capital Resources\b",
    r"\bResults of Operations\b", r"\bForward[- ]Looking\b"
]

SUPPORTED_EXTS = {".pdf", ".docx", ".txt", ".md", ".csv", ".html", ".htm", ".pptx"}

# ---------------------------
# Utils
# ---------------------------
def cosine_sim(a: np.ndarray, b: np.ndarray) -> float:
    na = np.linalg.norm(a); nb = np.linalg.norm(b)
    if na == 0 or nb == 0: return 0.0
    return float(np.dot(a, b) / (na * nb))

def split_into_sentences(text: str) -> List[str]:
    text = re.sub(r"\s+", " ", text).strip()
    sents = re.split(r"(?<=[^A-Z].[.?]) +(?=[A-Z(])", text)
    return [s.strip() for s in sents if s.strip()]

def looks_like_heading(line: str) -> bool:
    if len(line) > 180: return False
    if any(re.search(pat, line, flags=re.I) for pat in HARD_HEADINGS): return True
    if re.match(r"^[A-Z][A-Z0-9\s,&\-]{6,}$", line): return True
    return False

# ---------- Company line heuristics (avoid “as specified in its charter”) ----------
COMPANY_ENDINGS = (
    " INC", " INC.", " INCORPORATED", " CORPORATION", " CORP", " CORP.",
    " LTD", " LTD.", " LIMITED", " PLC", " LLC", " L.L.C", " NV", " N.V.",
    " AG", " SE", " SA", " S.A."
)
BAD_CUE_SUBSTRINGS = (
    "as specified in its charter",
    "exact name of registrant",
    "state or other jurisdiction",
    "commission file number",
    "irs employer identification",
    "principal executive offices",
)

def _clean_company(s: str) -> str:
    s = re.sub(r"\s+", " ", s).strip()
    return s.strip("•*-–—()[]:;. ")

def _looks_like_company_line(line: str) -> bool:
    L = _clean_company(line).upper()
    if len(L) < 3 or len(L) > 120:
        return False
    if any(b in L for b in BAD_CUE_SUBSTRINGS):
        return False
    if any(L.endswith(end) for end in COMPANY_ENDINGS):
        return True
    if re.match(r"^[A-Z0-9&.,'()\- ]{5,}$", L) and not re.search(r"\bREPORT\b|\bFORM\b|\bQUARTERLY\b|\bANNUAL\b", L):
        return True
    return False

def _guess_company_name_from_lines(lines: List[str]) -> Optional[str]:
    # Prefer the line above the SEC cue
    for i, ln in enumerate(lines):
        if "exact name of registrant" in ln.strip().lower() and i > 0:
            above = _clean_company(lines[i-1])
            if _looks_like_company_line(above):
                return above
    # Otherwise first plausible line near top
    for ln in lines[:40]:
        if _looks_like_company_line(ln):
            return _clean_company(ln)
    # Fallback anywhere
    for ln in lines:
        if any(_clean_company(ln).upper().endswith(end) for end in COMPANY_ENDINGS):
            return _clean_company(ln)
    return None

# ---------------------------
# Metadata
# ---------------------------
@dataclass
class FinancialMetadata:
    title: Optional[str]
    company: Optional[str]
    ticker: Optional[str]
    period: Optional[str]
    fiscal_year: Optional[str]
    source_path: Optional[str]
    pdf_meta: Dict[str, Any]

def _extract_financial_metadata_from_first_text(first_text: str, pdf_meta: Dict[str, Any], source_path: Optional[str]) -> FinancialMetadata:
    raw_lines = [ln for ln in (first_text.splitlines() if first_text else []) if ln.strip()]
    company = _guess_company_name_from_lines(raw_lines)
    block = re.sub(r"\s+", " ", first_text or "").strip()

    fiscal_year = None
    m_fy = re.search(r"\bFiscal\s+Year\s+(?:Ended|Ending)\b[:\s\-]*(.*?)(?:\.|,|$)", block, re.I)
    if m_fy: fiscal_year = m_fy.group(1).strip()
    else:
        m_fy = re.search(r"For the (?:fiscal )?year ended\s+([A-Za-z]+\s+\d{1,2},\s+\d{4})", block, re.I)
        if m_fy: fiscal_year = m_fy.group(1).strip()

    period = None
    m_period = re.search(r"Quarter(?:ly)? Report.*?ended\s+([A-Za-z]+\s+\d{1,2},\s+\d{4})", block, re.I)
    if m_period:
        period = f"Quarter ended {m_period.group(1)}"

    ticker = None
    m_ticker = re.search(r"(?:Trading\s+Symbol|Ticker)\s*[:\-]\s*([A-Z\.\-]{1,10})", block, re.I)
    if m_ticker:
        ticker = m_ticker.group(1).upper()

    title = (pdf_meta or {}).get("Title") or (company and f"{company} Financial Report") or "Financial Document"

    return FinancialMetadata(
        title=title, company=company, ticker=ticker, period=period,
        fiscal_year=fiscal_year, source_path=source_path, pdf_meta=pdf_meta or {}
    )

def extract_financial_metadata_from_text(text: str, filename: Optional[str] = None) -> FinancialMetadata:
    lines = [ln for ln in text.splitlines() if ln.strip()]
    company = _guess_company_name_from_lines(lines[:500])
    sample = re.sub(r"\s+", " ", "\n".join(lines[:500])).strip()

    fiscal_year = None
    m_fy = re.search(r"\bFiscal\s+Year\s+(?:Ended|Ending)\b[:\s\-]*(.*?)(?:\.|,|$)", sample, re.I)
    if m_fy: fiscal_year = m_fy.group(1).strip()
    else:
        m_fy = re.search(r"For the (?:fiscal )?year ended\s+([A-Za-z]+\s+\d{1,2},\s+\d{4})", sample, re.I)
        if m_fy: fiscal_year = m_fy.group(1).strip()

    period = None
    m_period = re.search(r"Quarter(?:ly)? Report.*?ended\s+([A-Za-z]+\s+\d{1,2},\s+\d{4})", sample, re.I)
    if m_period:
        period = f"Quarter ended {m_period.group(1)}"

    ticker = None
    m_ticker = re.search(r"(?:Trading\s+Symbol|Ticker)\s*[:\-]\s*([A-Z\.\-]{1,10})", sample, re.I)
    if m_ticker:
        ticker = m_ticker.group(1).upper()

    title = (company and f"{company} Financial Document") or (filename or "Financial Document")
    return FinancialMetadata(
        title=title, company=company, ticker=ticker, period=period,
        fiscal_year=fiscal_year, source_path=filename, pdf_meta={}
    )

# ---------------------------
# Extractors per file type (non-PDF)
# ---------------------------
def extract_text_from_docx_bytes(docx_bytes: bytes) -> str:
    if DocxDocument is None:
        raise RuntimeError("python-docx not installed. pip install python-docx")
    doc = DocxDocument(io.BytesIO(docx_bytes))
    parts = []
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text.strip())
    for tbl in doc.tables:
        for row in tbl.rows:
            row_txt = " | ".join(cell.text.strip() for cell in row.cells if cell.text)
            if row_txt:
                parts.append(row_txt)
    return "\n".join(parts).strip()

def extract_text_from_txt_bytes(txt_bytes: bytes, encoding="utf-8") -> str:
    return txt_bytes.decode(encoding, errors="ignore")

def extract_text_from_html_bytes(html_bytes: bytes, encoding="utf-8") -> str:
    if BeautifulSoup is None:
        raise RuntimeError("beautifulsoup4 not installed. pip install beautifulsoup4 lxml")
    html = html_bytes.decode(encoding, errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator=" ")
    return re.sub(r"\s+", " ", text).strip()

def extract_text_from_pptx_bytes(pptx_bytes: bytes) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx not installed. pip install python-pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
    return "\n".join(parts).strip()

def extract_text_any(file_bytes: bytes, filename: str) -> Tuple[str, str]:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        # handled elsewhere; caller will use PdfExtractor combined API
        raise RuntimeError("extract_text_any() should not be called for PDFs")
    if ext == ".docx":
        return extract_text_from_docx_bytes(file_bytes), "docx"
    if ext in (".txt", ".md", ".csv"):
        return extract_text_from_txt_bytes(file_bytes), "txt"
    if ext in (".html", ".htm"):
        return extract_text_from_html_bytes(file_bytes), "html"
    if ext == ".pptx":
        return extract_text_from_pptx_bytes(file_bytes), "pptx"
    raise ValueError(f"Unsupported file type: {ext}. Supported: {sorted(SUPPORTED_EXTS)}")

# ---------------------------
# Semantic chunker
# ---------------------------
class SemanticChunker:
    def __init__(self, embed_model_name: str):
        if not embed_model_name or not str(embed_model_name).strip():
            raise ValueError(
                "Embedding model name is empty. "
                "Set FIN_EMBED_MODEL or default to 'sentence-transformers/all-MiniLM-L6-v2'."
            )
        self.emb = HuggingFaceEmbeddings(model_name=embed_model_name)

    def embed(self, texts: List[str]) -> np.ndarray:
        vecs = self.emb.embed_documents(texts)
        return np.array(vecs, dtype=np.float32)

    def chunk(self, text: str) -> List[str]:
        sents = split_into_sentences(text)
        if not sents: return []

        chunks, cur, cur_vecs = [], [], []
        centroid = None

        for sent in sents:
            if looks_like_heading(sent) and sum(len(x) for x in cur) >= MIN_CHARS:
                chunks.append(" ".join(cur).strip())
                cur, cur_vecs, centroid = [], [], None

            cur.append(sent)

            if sum(len(x) for x in cur) >= MAX_CHARS:
                chunks.append(" ".join(cur).strip())
                cur, cur_vecs, centroid = [], [], None
                continue

            try:
                vec = self.embed([sent])[0]
                cur_vecs.append(vec)
                if centroid is None:
                    centroid = vec
                else:
                    sim = cosine_sim(centroid, vec)
                    if sim < 1.0:
                        centroid = (centroid * (len(cur_vecs)-1) + vec) / len(cur_vecs)
                    if sim < (1.0 - EMBED_BREAK_DELTA) and sum(len(x) for x in cur) > MIN_CHARS:
                        last = cur.pop()
                        chunks.append(" ".join(cur).strip())
                        cur, cur_vecs = [last], [vec]
                        centroid = vec
            except Exception:
                pass

        if cur:
            chunks.append(" ".join(cur).strip())
        return [c for c in chunks if c and len(c) > 50]

# ---------------------------
# Indexing & Retrieval
# ---------------------------
def build_or_load_index(docs: List[Document], emb_name: str, dirpath: Path) -> FAISS:
    dirpath.mkdir(parents=True, exist_ok=True)
    embeddings = HuggingFaceEmbeddings(model_name=emb_name)

    faiss_idx = dirpath / "index.faiss"
    faiss_pkl = dirpath / "index.pkl"
    if faiss_idx.exists() and faiss_pkl.exists():
        db = FAISS.load_local(folder_path=str(dirpath),
                              embeddings=embeddings,
                              allow_dangerous_deserialization=True)
    else:
        db = FAISS.from_documents(docs, embeddings)
        db.save_local(folder_path=str(dirpath))
    return db

@dataclass
class RAGSession:
    db: FAISS
    meta: FinancialMetadata
    docs: List[Document]

def make_documents_from_text(full_text: str, meta: FinancialMetadata) -> List[Document]:
    chunker = SemanticChunker(EMBED_MODEL)
    chunks = chunker.chunk(full_text)
    docs: List[Document] = []
    for i, ch in enumerate(chunks, start=1):
        docs.append(Document(
            page_content=ch,
            metadata={
                "chunk_id": i,
                "company": meta.company,
                "ticker": meta.ticker,
                "fiscal_year": meta.fiscal_year,
                "period": meta.period,
                "title": meta.title,
                "source": meta.source_path or "in-memory",
            }
        ))
    return docs

# ---------- Session builders ----------
def prepare_session_from_pdf_path(path: Union[str, Path]) -> RAGSession:
    extractor = PdfExtractor(SimplePdfValidator())
    full_text, first_text, pdf_meta = extractor.extract_text_and_header_from_path(str(path))
    meta = _extract_financial_metadata_from_first_text(first_text, pdf_meta, str(path))
    docs = make_documents_from_text(full_text, meta)
    db = build_or_load_index(docs, EMBED_MODEL, INDEX_DIR)
    return RAGSession(db=db, meta=meta, docs=docs)

def prepare_session_any(file_bytes: bytes, filename: str) -> RAGSession:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        extractor = PdfExtractor(SimplePdfValidator())
        buf = io.BytesIO(file_bytes); buf.name = filename
        full_text, first_text, pdf_meta = extractor.extract_text_and_header_from_file(buf)
        meta = _extract_financial_metadata_from_first_text(first_text, pdf_meta, filename)
    else:
        # non-PDF route
        text, _ftype = extract_text_any_nonpdf(file_bytes, filename)
        full_text = text
        meta = extract_financial_metadata_from_text(full_text, filename=filename)

    docs = make_documents_from_text(full_text, meta)
    db = build_or_load_index(docs, EMBED_MODEL, INDEX_DIR)
    return RAGSession(db=db, meta=meta, docs=docs)

# helper for non-pdf
def extract_text_any_nonpdf(file_bytes: bytes, filename: str) -> Tuple[str, str]:
    ext = Path(filename).suffix.lower()
    if ext == ".docx":
        return extract_text_from_docx_bytes(file_bytes), "docx"
    if ext in (".txt", ".md", ".csv"):
        return extract_text_from_txt_bytes(file_bytes), "txt"
    if ext in (".html", ".htm"):
        return extract_text_from_html_bytes(file_bytes), "html"
    if ext == ".pptx":
        return extract_text_from_pptx_bytes(file_bytes), "pptx"
    raise ValueError(f"Unsupported file type: {ext}. Supported: {sorted(SUPPORTED_EXTS)}")

# ---------------------------
# Retrieval APIs (no LLM)
# ---------------------------
def get_summary_context(session: RAGSession, max_chunks: int = 8) -> Tuple[str, FinancialMetadata]:
    ctx = "\n\n".join(d.page_content for d in session.docs[:max_chunks])
    return ctx, session.meta

def retrieve(session: RAGSession, query: str, k: int = TOP_K) -> Dict[str, Any]:
    retriever = session.db.as_retriever(search_kwargs={"k": k})
    docs = retriever.get_relevant_documents(query)
    context = "\n\n".join(d.page_content for d in docs)
    sources = []
    for d in docs:
        sources.append({
            "source": d.metadata.get("source", "unknown"),
            "chunk_id": d.metadata.get("chunk_id"),
            "company": d.metadata.get("company"),
        })
    return {"context": context, "sources": sources}